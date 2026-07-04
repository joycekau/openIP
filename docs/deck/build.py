# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image
import os, math

os.chdir(os.path.dirname(os.path.abspath(__file__)))

# ---------- palette ----------
BG="080510"; PANEL="150E26"; PANEL2="1B1232"; TEXT="F4EFFC"; MUTED="A99FC0"; DIM="7A6F96"
CYAN="22D3EE"; VIOLET="8B5CF6"; PURPLE="A855F7"; MAGENTA="EC4899"; PINK="F472B6"; CORAL="FB7185"; LINE="2A2140"
GRAD=[CYAN,VIOLET,MAGENTA,PINK,CORAL]
DISP="Space Grotesk"          # latin display / brand / numbers

# globals reset per language
LANG="zh"; BODY="Microsoft YaHei"; prs=None; TOTAL=16
def t(zh, en): return zh if LANG=="zh" else en
def C(h): return RGBColor.from_string(h)

def slide(bg="bg_content.png"):
    s=prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_picture(bg,0,0,width=prs.slide_width,height=prs.slide_height); return s

def shadow(shape, blur=20, dist=7, angle=90, color="000000", alpha=58):
    sp=shape._element.spPr; ex=sp.find(qn('a:effectLst'))
    if ex is not None: sp.remove(ex)
    sp.append(parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{int(blur*12700)}" dist="{int(dist*12700)}" dir="{int(angle*60000)}" rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{int(alpha*1000)}"/></a:srgbClr></a:outerShdw></a:effectLst>'))

def grad_fill(shape, stops, angle=45):
    sp=shape._element.spPr
    for tag in ('a:solidFill','a:noFill','a:gradFill','a:blipFill','a:pattFill'):
        e=sp.find(qn(tag))
        if e is not None: sp.remove(e)
    gs="".join(f'<a:gs pos="{int(p*1000)}"><a:srgbClr val="{c}"/></a:gs>' for p,c in stops)
    xml=('<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">'
         f'<a:gsLst>{gs}</a:gsLst><a:lin ang="{int(angle*60000)}" scaled="1"/></a:gradFill>')
    ln=sp.find(qn('a:ln'))
    if ln is not None: ln.addprevious(parse_xml(xml))
    else: sp.append(parse_xml(xml))

def card(s,x,y,w,h,fill=PANEL,radius=0.085,line=LINE,line_w=1.0,sh=True):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.adjustments[0]=radius; shp.fill.solid(); shp.fill.fore_color.rgb=C(fill)
    if line: shp.line.color.rgb=C(line); shp.line.width=Pt(line_w)
    else: shp.line.fill.background()
    if sh: shadow(shp)
    return shp

def dot(s,x,y,d,fill,glow=False):
    shp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb=C(fill); shp.line.fill.background()
    if glow: shadow(shp,blur=22,dist=0,alpha=45,color=fill)
    return shp

def txt(s,runs,x,y,w,h,size=16,color=TEXT,bold=False,font=None,
        align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spc=None,leading=None,sa=None,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    if isinstance(runs,str): runs=[[(runs,color,bold,size,font)]]
    elif runs and isinstance(runs[0],tuple): runs=[runs]
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align
        if leading: p.line_spacing=leading
        if sa is not None: p.space_after=Pt(sa)
        p.space_before=Pt(0)
        for (tt,col,bd,sz,fn) in para:
            fn_use=fn if fn else BODY
            r=p.add_run(); r.text=tt; r.font.size=Pt(sz); r.font.bold=bd
            r.font.color.rgb=C(col); r.font.name=fn_use
            rPr=r._r.get_or_add_rPr(); ea=rPr.find(qn('a:ea'))
            if ea is None: ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
            ea.set('typeface', BODY)
            if spc is not None: rPr.set('spc', str(int(spc*100)))
    return tb

def R(tt,col=TEXT,bd=False,sz=16,fn=None): return (tt,col,bd,sz,fn)
def eyebrow(s,tt,x,y,color=MAGENTA,w=9): txt(s,[[R(tt,color,True,12.5,DISP)]],x,y,w,0.3,spc=2.4)
def title(s,tt,x=0.7,y=0.78,size=33,w=11.6,color=TEXT): txt(s,[[R(tt,color,True,size)]],x,y,w,0.9)
def notes(s,zh,en): s.notes_slide.notes_text_frame.text = zh if LANG=="zh" else en

mw,mh=Image.open("onee.png").size; ASP=mw/mh
def mascot(s,x,y,w): s.shapes.add_picture("onee.png",Inches(x),Inches(y),width=Inches(w))

def logo_lockup(s,x,y,h,tsize,text="neIP.io",gap=0.06,w=7.5):
    s.shapes.add_picture("logo_mark.png",Inches(x),Inches(y),height=Inches(h))
    runs=[]; n=len(text)
    for i,ch in enumerate(text):
        runs.append(R(ch,GRAD[min(int(i*len(GRAD)/max(n,1)),len(GRAD)-1)],True,tsize,DISP))
    txt(s,[runs],x+h+gap,y,w,h,anchor=MSO_ANCHOR.MIDDLE)

def feature_pills(s,x,y,pills,pw=1.92,ph=0.6,gap=0.2):
    cxp=x
    for ic,lab,col in pills:
        card(s,cxp,y,pw,ph,fill=PANEL2,radius=0.32,line=col,line_w=1.1,sh=False)
        txt(s,[[R(ic,col,True,15,DISP)]],cxp+0.22,y,0.5,ph,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,[[R(lab,TEXT,True,13.5)]],cxp+0.6,y,pw-0.72,ph,anchor=MSO_ANCHOR.MIDDLE)
        cxp+=pw+gap

def pageno(s):
    idx=len(prs.slides._sldIdLst)
    s.shapes.add_picture("logo_mark.png",Inches(0.7),Inches(7.0),height=Inches(0.33))
    txt(s,[[R("oneIP.io",MUTED,True,10,DISP),R(t("  ·  融资计划书","  ·  Pitch Deck"),DIM,False,10)]],
        1.12,7.0,5,0.33,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,[[R(f"{idx:02d} / {TOTAL}",DIM,False,10,DISP)]],11.4,7.0,1.3,0.33,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)

# =================================================================== BUILD
def build():
    # ---- S1 COVER ----
    s=slide("bg_cover.png")
    eyebrow(s,t("网红 IP 发币社交平台 · 融资计划书 2026","INFLUENCER IP LAUNCH & SOCIAL PLATFORM · PITCH 2026"),0.75,0.62,CYAN,w=12)
    logo_lockup(s,0.7,1.12,1.35,66)
    txt(s,[[R("OneIP",CYAN,True,33,DISP),R(" + ",DIM,False,33,DISP),R("OPC",PINK,True,33,DISP)]],0.78,2.5,8,0.65)
    txt(s,[[R("OPC = One Personal Company",TEXT,True,14,DISP),
            R(t("　一人即一家公司","  ·  you are the company"),MUTED,False,14)]],0.78,3.22,8.7,0.38)
    if LANG=="zh":
        txt(s,[[R("品牌 IP 一键发币",PINK,True,18),R(" · ",DIM,False,18),R("链接粉丝",CYAN,True,18),
                R(" · ",DIM,False,18),R("价值共创",VIOLET,True,18)]],0.78,3.78,8.7,0.45)
    else:
        txt(s,[[R("Launch Brand IP",PINK,True,18,DISP),R(" · ",DIM,False,18,DISP),R("Connect Fans",CYAN,True,18,DISP),
                R(" · ",DIM,False,18,DISP),R("Co-create Value",VIOLET,True,18,DISP)]],0.78,3.78,8.9,0.45)
    feature_pills(s,0.78,4.68,[("◆",t("发币","Launch"),CYAN),("◎",t("社交","Social"),VIOLET),
                               ("✦",t("共创","Co-create"),MAGENTA),("★",t("激励","Rewards"),CORAL)])
    chips=[(t("融资","Raise"),t("200 万美金","US$2M"),CYAN),
           (t("让出原始股","Equity"),"20%",PINK),
           (t("投后估值","Post-money"),t("1,000 万美金","US$10M"),VIOLET)]
    cx=0.78
    for lab,val,col in chips:
        card(s,cx,5.95,2.55,0.95,fill=PANEL2,radius=0.22,line=LINE,line_w=1)
        txt(s,[[R(lab,MUTED,False,11)]],cx+0.28,6.12,2.05,0.3)
        txt(s,[[R(val,col,True,19,DISP)]],cx+0.28,6.42,2.05,0.4)
        cx+=2.83
    mascot(s,9.5,1.5,3.5)

    # ---- S2 MEGATREND (industry forecast) ----
    s=slide("bg_section.png")
    eyebrow(s,t("MEGATREND · 结构性趋势（行业预测）","MEGATREND · INDUSTRY FORECAST"),0.7,0.58,CYAN,w=12)
    title(s,t("AI 越强，个人 IP 越值钱","The Stronger AI Gets, the More Personal IP Is Worth"),y=0.9,size=28)
    txt(s,[[R(t("机构预测（IDC · Stanford HAI · Goldman Sachs · WEF · McKinsey）：AI 经济狂飙、岗位重构，个人品牌成为出口。",
        "Institutional forecasts (IDC · Stanford HAI · Goldman Sachs · WEF · McKinsey): the AI economy soars, jobs restructure, personal brand becomes the exit."),MUTED,False,12.5)]],0.72,1.54,12.1,0.4)
    yrs=["2027","2028","2029","2030"]; yx=[5.0,6.85,8.7,10.55]; yw=1.75
    txt(s,[[R(t("指标 · 行业预测","Metric · forecast"),DIM,True,10.5)]],0.85,1.98,3.9,0.3,spc=0.6)
    for x,yy in zip(yx,yrs): txt(s,[[R(yy,DIM,True,11,DISP)]],x,1.98,yw,0.3,align=PP_ALIGN.CENTER)
    trends=[(t("AI 全球经济价值","Global AI economic value"),["$3.2T","$3.8T","$4.4T","$4.9T"],CYAN),
            (t("AI 市场年成长 (CAGR)","AI market CAGR"),["25%","27%","28%","30%"],VIOLET),
            (t("企业员工减少 (预测)","Workforce cut (forecast)"),["5–10%","10–15%","15–20%","20–30%"],MAGENTA),
            (t("AI 可自动化岗位比例","Automatable jobs"),["20%","25%","30%","35%+"],CORAL),
            (t("创作者经济规模","Creator economy"),["$400B","$440B","$470B","$500B+"],PINK),
            (t("AI 技能需求","AI-skill demand"),[t("高","High"),t("很高","Higher"),t("标配","Standard"),t("必备","Must-have")],PURPLE)]
    ry,rh,gap=2.26,0.54,0.05
    for lab,vals,col in trends:
        card(s,0.6,ry,12.13,rh,fill=PANEL,radius=0.16,sh=False)
        dot(s,0.85,ry+rh/2-0.05,0.11,col)
        txt(s,[[R(lab,TEXT,True,11.5)]],1.06,ry,3.75,rh,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
        for i,(x,v) in enumerate(zip(yx,vals)):
            last=(i==3)
            txt(s,[[R(v,(col if last else MUTED),last,(14 if last else 12),DISP)]],x,ry,yw,rh,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        ry+=rh+gap
    card(s,0.6,5.85,12.13,0.72,fill=PANEL2,radius=0.16,line=VIOLET,line_w=1.3)
    txt(s,[[R(t("专业人士转型自媒体（Creator Economy）：","Professionals pivot to the Creator Economy: "),CYAN,True,13),
            R(t("医生 · 律师 · 会计 · 老师 · 地产 · 保险 纷纷经营个人 IP —— 正是 OneIP 的用户。",
                "doctors · lawyers · accountants · teachers · realtors · agents all build personal IP — OneIP's users."),TEXT,False,11.5)]],
        0.95,5.85,11.5,0.72,anchor=MSO_ANCHOR.MIDDLE,leading=1.08)
    txt(s,[[R(t("来源：IDC / Stanford AI Index / Goldman Sachs / WEF / McKinsey（2024–2026）；「企业员工减少」为综合机构预测区间，非官方统计。",
        "Sources: IDC / Stanford AI Index / Goldman Sachs / WEF / McKinsey (2024–2026). \"Workforce cut\" is a blended industry-forecast range, not official data."),DIM,False,8.5)]],0.72,6.64,12.2,0.26)
    notes(s,
        "数据来源（行业预测 / 已发生趋势，2024–2026）：\n"
        "• AI 全球经济价值：IDC 预测 AI 到 2030 为全球经济贡献约 19.9 万亿美元；本页取核心直接贡献口径 $3.2T→$4.9T，年增约 25–30% CAGR。\n"
        "  Axios/IDC: https://www.axios.com/2024/09/17/ai-global-economy-idc-2030\n"
        "• 企业 AI 使用率 / AI Agent 普及：Stanford AI Index (HAI)。https://hai.stanford.edu/ai-index\n"
        "• 企业员工减少 / AI 可自动化岗位：McKinsey、WEF 估 20–35% 工作可被自动化；Goldman Sachs 估约 3 亿全职岗位受影响。为综合机构预测区间，非官方统计。\n"
        "  WEF: https://www.businessinsider.com/wef-sees-4-ai-futures-for-jobs-by-2030-only-one-limits-disruption-2026-1\n"
        "  Goldman/Investopedia: https://www.investopedia.com/if-you-lose-your-job-to-a-i-the-setback-might-be-permanent-11945756\n"
        "• 毕业生 / 初级岗位：Reuters 报道中国企业减少毕业生招聘、AI 取代部分初级岗位；Business Insider 指出各产业差异大。\n"
        "  Reuters: https://www.reuters.com/business/world-at-work/china-inc-deploys-quiet-layoffs-beijing-promotes-ai-adoption-2026-06-10/\n"
        "  BI: https://www.businessinsider.com/the-ai-layoffs-story-just-got-more-complicated-2026-6\n"
        "• 创作者经济：Goldman Sachs 估 2027 接近 $500B；越来越多专业人士（医生/律师/会计/老师/地产/保险）经营个人 IP，AI 大幅降低制作成本。\n"
        "  https://en.wikipedia.org/wiki/Creator_economy\n"
        "注：数字为区间/预测，用于说明结构性方向，非业绩承诺。",
        "Data sources (industry forecasts / current trends, 2024–2026):\n"
        "• Global AI economic value: IDC forecasts AI to add ~US$19.9T to the global economy by 2030; this page uses a core direct-contribution view $3.2T→$4.9T, ~25–30% CAGR.\n"
        "  Axios/IDC: https://www.axios.com/2024/09/17/ai-global-economy-idc-2030\n"
        "• Enterprise AI adoption / AI agents: Stanford AI Index (HAI). https://hai.stanford.edu/ai-index\n"
        "• Workforce cut / automatable jobs: McKinsey & WEF estimate 20–35% of work is automatable; Goldman Sachs estimates ~300M full-time jobs affected. Blended industry-forecast range, not official data.\n"
        "  WEF: https://www.businessinsider.com/wef-sees-4-ai-futures-for-jobs-by-2030-only-one-limits-disruption-2026-1\n"
        "  Goldman/Investopedia: https://www.investopedia.com/if-you-lose-your-job-to-a-i-the-setback-might-be-permanent-11945756\n"
        "• Graduates / entry-level: Reuters reports Chinese firms cutting graduate hiring as AI replaces some junior roles; BI notes wide sector variation.\n"
        "  Reuters: https://www.reuters.com/business/world-at-work/china-inc-deploys-quiet-layoffs-beijing-promotes-ai-adoption-2026-06-10/\n"
        "  BI: https://www.businessinsider.com/the-ai-layoffs-story-just-got-more-complicated-2026-6\n"
        "• Creator economy: Goldman Sachs estimates ~US$500B by 2027; more professionals (doctors/lawyers/accountants/teachers/realtors/agents) build personal IP as AI slashes production cost.\n"
        "  https://en.wikipedia.org/wiki/Creator_economy\n"
        "Note: figures are ranges/forecasts illustrating structural direction, not performance commitments.")
    pageno(s)

    # ---- S3 AI ERA ----
    s=slide("bg_content.png")
    eyebrow(s,t("THE AI ERA · 时代背景","THE AI ERA"),0.7,0.6,CYAN)
    title(s,t("AI 正在重写每一条产业线","AI Is Rewriting Every Industry"),y=0.95)
    txt(s,[[R(t("AI 的进化速度超出人类想象，实体产业正被加速替代——而价值，正在向「无法被复制的东西」迁移。",
        "AI is evolving faster than anyone imagined. As it replaces real-world work, value migrates to what cannot be copied."),MUTED,False,15.5)]],
        0.72,1.7,11.7,0.6,leading=1.25)
    cards2=[("⚡",t("进化超预期","Beyond Expectation"),t("AI 能力以月为单位迭代，远快于任何一次工业革命的节奏。","AI capability iterates monthly — faster than any industrial revolution."),CYAN),
            ("🏭",t("产业全面引进","Adopted Everywhere"),t("客服、设计、文案、编程、客拓——岗位被一一接管。","Support, design, copy, coding, sales — roles taken over one by one."),VIOLET),
            ("📉",t("事业被取代","Jobs Replaced"),t("越来越多人因 AI 失去当前职业，技能护城河快速蒸发。","More people lose their jobs to AI; skill moats evaporate fast."),MAGENTA),
            ("❓",t("价值在迁移","Value Migrates"),t("可被复制的劳动趋于零成本；稀缺性成为唯一的溢价来源。","Copyable labor trends to zero cost; scarcity becomes the only premium."),CORAL)]
    x0,y0,cw,ch,gap=0.72,2.55,2.92,3.65,0.18
    for i,(ic,h,b,col) in enumerate(cards2):
        x=x0+i*(cw+gap)
        card(s,x,y0,cw,ch,fill=PANEL,radius=0.07)
        dot(s,x+0.32,y0+0.34,0.62,PANEL2)
        txt(s,[[R(ic,col,True,22,DISP)]],x+0.32,y0+0.40,0.62,0.5,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,[[R(h,TEXT,True,17.5)]],x+0.32,y0+1.2,cw-0.6,0.5)
        txt(s,[[R(b,MUTED,False,12.5)]],x+0.32,y0+1.72,cw-0.6,1.8,leading=1.3)
    pageno(s)

    # ---- S3 INSIGHT ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE INSIGHT · 不可被取代","THE INSIGHT"),0.7,0.62,CYAN)
    title(s,t("唯独「人」和「粉丝」，AI 无法取代","Only People & Fans Can't Be Replaced"),y=0.97)
    left=[(t("个人 IP / 自媒体","Personal IP / Creators"),t("AI 只能辅助创作，无法成为「那个人」。个性、故事、信任——天然独一无二。","AI can assist, but can't BE that person. Personality, story, trust — inherently unique."),VIOLET),
          (t("粉丝数据 / 粉丝情绪","Fan Data / Fan Emotion"),t("真实的关注与情感连接无法被生成。AI 是放大器，不是替代品。","Real attention and emotion can't be generated. AI amplifies — it doesn't replace."),MAGENTA)]
    ly=2.05
    for h,b,col in left:
        card(s,0.72,ly,6.5,1.95,fill=PANEL,radius=0.08)
        dot(s,1.02,ly+0.32,0.22,col,glow=True)
        txt(s,[[R(h,TEXT,True,19)]],1.45,ly+0.28,5.5,0.5)
        txt(s,[[R(b,MUTED,False,13)]],1.02,ly+0.95,5.95,0.9,leading=1.3)
        ly+=2.15
    card(s,7.55,2.05,5.05,4.05,fill=PANEL2,radius=0.08,line=VIOLET,line_w=1.2)
    if LANG=="zh":
        txt(s,[[R("一个 IP，一个品牌",TEXT,True,22)],[R("粉丝情绪，产生了价值",TEXT,True,22)]],7.95,2.55,4.3,1.2,leading=1.2,sa=2)
    else:
        txt(s,[[R("One IP, one brand —",TEXT,True,21)],[R("fan emotion creates value.",TEXT,True,21)]],7.95,2.55,4.4,1.2,leading=1.2,sa=2)
    txt(s,[[R(t("独一无二","One of a Kind"),CYAN,True,17)]],7.95,3.95,4.3,0.4)
    logo_lockup(s,7.95,4.45,0.92,38,w=4)
    txt(s,[[R(t("把这份独一无二，变成可定价、可流通、可变现的资产。",
        "Turn that uniqueness into a priceable, tradable, monetizable asset."),MUTED,False,13)]],7.95,5.55,4.4,0.7,leading=1.3)
    pageno(s)

    # ---- S4 DUAL ENGINE ----
    s=slide("bg_content.png")
    eyebrow(s,t("THE SOLUTION · 双引擎","THE SOLUTION"),0.7,0.6,CYAN)
    title(s,t("OneIP：一个 IP 经济的双引擎闭环","OneIP: A Dual-Engine IP Economy"),y=0.95)
    eng=[("oneIP.ai",t("造 IP · 供给端","Create IP · Supply"),
          t("用 AI 把任何创作者 / KOL / 品牌，批量生产成可分发、可授权的数字 IP 资产。",
            "Use AI to turn any creator / KOL / brand into distributable, licensable digital IP assets at scale."),CYAN,
          [t("AI 数字分身 / IP 形象","AI avatar / IP persona"),t("内容批量生产（图文 / 视频）","Content at scale (text / video)"),t("标准化「IP 资产包」授权","Licensable IP asset packs")]),
         ("oneIP.io",t("卖 IP · 流通端","Trade IP · Liquidity"),
          t("基于 Solana 的一键发行与交易终端（pump.fun 式），让 IP 资产链上定价、流通、变现。",
            "A Solana launch & trading terminal (pump.fun-style) that prices, circulates and monetizes IP on-chain."),PINK,
          [t("一键发行 IP / KOL 代币","1-click IP / KOL token launch"),t("pump.fun 式交易终端","pump.fun-style trading terminal"),t("手续费 + 流动性变现","Fees + liquidity monetization")])]
    ex=0.72
    for name,sub,desc,col,items in eng:
        card(s,ex,2.05,5.62,4.6,fill=PANEL,radius=0.07)
        txt(s,[[R(name,col,True,30,DISP)]],ex+0.45,2.35,5.0,0.6)
        txt(s,[[R(sub,TEXT,True,15)]],ex+0.45,3.0,5.0,0.4)
        txt(s,[[R(desc,MUTED,False,13)]],ex+0.45,3.5,5.0,1.15,leading=1.32)
        iy=4.8
        for it in items:
            dot(s,ex+0.48,iy+0.07,0.14,col)
            txt(s,[[R(it,TEXT,False,13)]],ex+0.78,iy-0.04,5.0-1.2,0.4)
            iy+=0.52
        ex+=5.94
    txt(s,[[R("↔",VIOLET,True,30,DISP)]],6.34,4.0,0.66,0.6,align=PP_ALIGN.CENTER)
    mascot(s,11.75,0.45,1.4)
    pageno(s)

    # ---- S5 MODEL (ai + io merged) ----
    s=slide("bg_content.png")
    eyebrow(s,t("BUSINESS MODEL · 双引擎","BUSINESS MODEL · DUAL ENGINE"),0.7,0.58,CYAN,w=12)
    title(s,t("双引擎商业模式：造 IP × 卖 IP","Dual-Engine Model: Create IP × Trade IP"),y=0.9,size=29)
    txt(s,[[R(t("oneIP.ai 造 IP 供给，oneIP.io 定价与变现——一套有真实资产支撑的 IP 经济闭环。",
        "oneIP.ai supplies IP, oneIP.io prices & monetizes it — one real-asset-backed IP economy."),MUTED,False,14)]],0.72,1.58,12,0.4)
    engines=[("oneIP.ai",t("造 IP · 供给端","Create IP · Supply"),CYAN,
              [(t("产品","Product"),t("AI 数字分身 · 内容批量生产 · IP 资产包","AI avatar · content at scale · IP packs")),
               (t("客户","Customers"),t("KOL · 艺人 · 品牌 · MCN","KOLs · talent · brands · MCNs")),
               (t("收入","Revenue"),t("AI 订阅(SaaS) · 授权抽成 · 企业定制","AI subs · licensing · custom IP")),
               (t("护城河","Moat"),t("IP 资产库飞轮 · 一致性模型微调","IP-library flywheel · consistency models"))]),
             ("oneIP.io",t("卖 IP · 流通端","Trade IP · Liquidity"),PINK,
              [(t("产品","Product"),t("一键发币 · pump.fun 式终端 · 钱包行情","1-click launch · pump.fun-style terminal")),
               (t("用户","Users"),t("创作者(发行) · 粉丝散户(交易)","Creators issue · fans/retail trade")),
               (t("收入","Revenue"),t("交易手续费(核心) · 发行费 · 增值订阅","Trading fees (core) · launch · subs")),
               (t("护城河","Moat"),t("交易量→流动性网络效应 · 用户黏性","Volume→liquidity network · stickiness"))])]
    ex=0.6
    for name,tag,col,blocks in engines:
        w=6.0
        card(s,ex,2.1,w,4.3,fill=PANEL,radius=0.07)
        txt(s,[[R(name,col,True,28,DISP)]],ex+0.4,2.3,w-0.8,0.55)
        txt(s,[[R(tag,TEXT,True,13.5)]],ex+0.4,2.85,w-0.8,0.35)
        by=3.35
        for blab,btext in blocks:
            dot(s,ex+0.42,by+0.08,0.12,col)
            txt(s,[[R(blab,col,True,12.5)]],ex+0.66,by-0.02,1.5,0.32)
            txt(s,[[R(btext,MUTED,False,11.5)]],ex+0.66,by+0.3,w-1.05,0.42,leading=1.12)
            by+=0.75
        ex+=6.13
    txt(s,[[R(t("风控铁律：","Iron rule: "),CORAL,True,12),
            R(t("不充值直换 · 不拿充值保收益 · 变现只走外部市场——守防庞氏。",
                "No top-up swaps · no guaranteed returns · external cash-out only — anti-Ponzi."),TEXT,False,11.5)]],0.72,6.5,12.2,0.35)
    pageno(s)

    # ---- S7 FLYWHEEL ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE FLYWHEEL · 商业闭环","THE FLYWHEEL"),0.7,0.62,CYAN)
    title(s,t("AI 供给 × 链上流通 = 自增长飞轮","AI Supply × On-chain Liquidity = Self-growing Flywheel"),y=0.97,size=30)
    txt(s,[[R(t("别人只做发币（投机、易归零）；我们用 AI 端保证真实供给，用链上端提供流动性与变现。",
        "Others only launch tokens (speculative, prone to zero). We secure real supply via AI and liquidity via chain."),MUTED,False,15)]],0.72,1.72,11.8,0.5,leading=1.25)
    cxp,cyp=6.666,4.45
    nodes=[(t("AI 持续生产 IP 资产","AI keeps producing IP assets"),CYAN,0.0),
           (t("链上发行 · 提供流动性","On-chain launch · liquidity"),PINK,90.0),
           (t("交易变现 · 反哺创作者","Trading monetizes · rewards creators"),MAGENTA,180.0),
           (t("更多 IP / 品牌入驻","More IPs / brands onboard"),VIOLET,270.0)]
    RX,RY,nw,nh=3.7,1.75,3.5,0.95
    hub=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cxp-1.05),Inches(cyp-1.05),Inches(2.1),Inches(2.1))
    grad_fill(hub,[(0,CYAN),(45,VIOLET),(100,MAGENTA)],angle=60); hub.line.fill.background(); shadow(hub,blur=26,dist=0,alpha=40,color=VIOLET)
    txt(s,[[R("OneIP","FFFFFF",True,21,DISP)],[R(t("飞轮","Flywheel"),"FFFFFF",True,13)]],cxp-1.05,cyp-0.62,2.1,1.2,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
    for label,col,ang in nodes:
        a=math.radians(ang); nx=cxp+RX*math.cos(a)-nw/2; ny=cyp+RY*math.sin(a)-nh/2
        card(s,nx,ny,nw,nh,fill=PANEL2,radius=0.2,line=col,line_w=1.3)
        txt(s,[[R(label,TEXT,True,13.5)]],nx+0.2,ny,nw-0.4,nh,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.05)
    pageno(s)

    # ---- S8 VALUE MATRIX (merged) ----
    s=slide("bg_content.png")
    eyebrow(s,t("WHO WINS · 合作方价值","WHO WINS · VALUE FOR PARTNERS"),0.7,0.58,CYAN,w=12)
    title(s,t("合作方价值：人人都能在 OneIP 赚到钱","Value for Everyone: Every Side Wins on OneIP"),y=0.9,size=29)
    txt(s,[[R(t("无论你是品牌、供应商，还是普通人、网红、品牌 IP，都能在 OneIP 找到自己的赚钱方式。",
        "Brand, supplier, everyday creator, KOL or IP — each side has its own way to earn on OneIP."),MUTED,False,14)]],0.72,1.62,12,0.4)
    dimheads=[t("为什么与我们合作","Why partner with us"),t("在 OneIP 赚什么","What you earn on OneIP"),t("OneIP 帮你什么","How OneIP helps")]
    colx=[3.05,6.29,9.53]; colw=3.16
    for cxh,hh in zip(colx,dimheads):
        txt(s,[[R(hh,CYAN,True,12.5)]],cxh,2.02,colw,0.3)
    audiences=[
        (t("品牌 & 供应商","Brands & Suppliers"),t("商业侧","Commerce side"),MAGENTA,
         [[t("直达 KOL 与高粘性粉丝","Reach KOLs & loyal fans"),t("粉丝即消费者，转化更高","Fans are buyers — higher conversion"),t("一站对接海量 IP","One platform, countless IPs")],
          [t("粉丝社区直接带货","Sell into fan communities"),t("IP 联名 / 授权分成","IP collab / licensing share"),t("早期发行资产增值","Upside from early launches")],
          [t("IP + 粉丝数据精准匹配","Precise IP + fan-data match"),t("链上自动分账、透明","On-chain auto split, transparent"),t("守防庞氏合规环境","Anti-Ponzi, compliant venue")]]),
        (t("普通人 · 网红 · 品牌 IP","Creators · KOLs · IP"),t("创作侧","Creator side"),CYAN,
         [[t("影响力首次可资产化","Influence becomes an asset"),t("零门槛发币，无需懂技术","Zero-barrier token launch"),t("AI 辅助持续产内容","AI-assisted content engine")],
          [t("发币募资 + 持续版税","Launch raise + ongoing royalties"),t("手续费 · 打赏 · 会员订阅","Fees · tips · memberships"),t("品牌合作 / 带货分成","Brand deals / sales share")],
          [t("一键发币 + 交易终端","1-click launch + terminal"),t("AI 造 IP：分身 / 内容","AI-built IP: avatar & content"),t("社交激励 + 流量资源","Social tools + traffic & brands")]])]
    ry=2.5
    for aud,tag,acol,cells in audiences:
        card(s,0.6,ry,12.13,1.98,fill=PANEL,radius=0.06)
        dot(s,0.92,ry+0.34,0.2,acol,glow=True)
        txt(s,[[R(aud,TEXT,True,15.5)]],1.0,ry+0.62,2.0,0.9,leading=1.08)
        txt(s,[[R(tag,acol,True,11)]],0.95,ry+1.45,2.0,0.3)
        for ci,cx in enumerate(colx):
            iy=ry+0.22
            for b in cells[ci]:
                dot(s,cx,iy+0.1,0.1,acol)
                txt(s,[[R(b,MUTED,False,11)]],cx+0.24,iy,colw-0.3,0.55,leading=1.08)
                iy+=0.56
        ry+=2.13
    pageno(s)

    # ---- S9 LANDSCAPE (GMGN removed) ----
    s=slide("bg_content.png")
    eyebrow(s,t("LANDSCAPE · 竞争格局","LANDSCAPE"),0.7,0.58,CYAN)
    title(s,t("为什么是 OneIP：一个平台补齐所有人的短板","Why OneIP: One Platform Fixes Everyone's Gaps"),y=0.9,size=29)
    txt(s,[[R(t("发币、交易、社交、品牌——各类玩家各有所长却各有缺口；OneIP 用一个平台整合，并守住安全底线。",
        "Launch, trading, social, brands — each player has strengths and gaps. OneIP unifies them, safely."),MUTED,False,14)]],0.72,1.58,12,0.4)
    HEADS=[t("对象 / 平台","Player"),t("核心功能","Core function"),t("特点","Trait"),t("安全 · 风控","Safety"),t("痛点","Pain point"),t("亮点","Edge")]
    hx0=[0.82,2.35,4.35,6.2,8.0,10.4]
    for cxh,hh in zip(hx0,HEADS): txt(s,[[R(hh,DIM,True,10.5)]],cxh,2.05,2.2,0.3,spc=0.6)
    CELLCOLS=[(2.35,1.95),(4.35,1.8),(6.2,1.75),(8.0,2.35),(10.4,2.3)]
    ROWS=[("pump.fun",[t("一键发 meme 币","1-click meme launch"),t("纯链上 · 无门槛","Pure on-chain · open"),t("无审核 · 易 Rug","No vetting · rug-prone"),t("投机归零 · 无沉淀","Speculative · no retention"),t("冷启动快","Fast cold-start")],CORAL),
          (t("现有社交平台","Social platforms"),[t("内容分发 / 连接","Content & connection"),t("流量中心化","Centralized reach"),t("平台合规","Platform-compliant"),t("难变现 · 高抽成","Hard to monetize · high cut"),t("粉丝基础大","Large fan base")],VIOLET),
          (t("传统品牌 / IP","Brands / IP"),[t("授权 / 联名变现","Licensing / collab"),t("中心化授权","Centralized licensing"),t("合规成熟","Mature compliance"),t("链条慢 · 零参与","Slow · no fan stake"),t("品牌信任高","High trust")],CYAN),
          ("OneIP",[t("发币+交易+社交+AI","Launch+trade+social+AI"),t("IP 资产化 · 粉丝共创","IP-ization · co-creation"),t("守防庞氏 · 外部变现","Anti-Ponzi · ext. cash-out"),t("补齐全部短板","Fixes every gap"),t("真实价值闭环","Real-value loop")],None)]
    ry,rh,gap=2.45,0.88,0.12
    for name,cells,col in ROWS:
        hot=(col is None)
        if hot:
            card(s,0.6,ry,12.13,rh,fill=PANEL2,radius=0.1,line=VIOLET,line_w=1.5)
            s.shapes.add_picture("logo_mark.png",Inches(0.8),Inches(ry+rh/2-0.19),height=Inches(0.38))
            txt(s,[[R(name,TEXT,True,12.5,DISP)]],1.26,ry,1.1,rh,anchor=MSO_ANCHOR.MIDDLE)
        else:
            card(s,0.6,ry,12.13,rh,fill=PANEL,radius=0.1,line=LINE,line_w=1.0,sh=False)
            dot(s,0.84,ry+rh/2-0.055,0.11,col)
            txt(s,[[R(name,TEXT,True,11,fn=(DISP if name=="pump.fun" else None))]],1.06,ry,1.25,rh,anchor=MSO_ANCHOR.MIDDLE,leading=1.02)
        for i,(cx,cw) in enumerate(CELLCOLS):
            cc=([TEXT,CYAN,"34D399","34D399",PINK] if hot else ["C9C0E0","A99FC0","A99FC0","D7A6B5","C9C0E0"])[i]
            txt(s,[[R(cells[i],cc,(hot and i==4),10)]],cx,ry,cw,rh,anchor=MSO_ANCHOR.MIDDLE,leading=1.1)
        ry+=rh+gap
    pageno(s)

    # ---- S10 MARKET TAM ----
    s=slide("bg_content.png")
    eyebrow(s,t("MARKET · 市场规模","MARKET"),0.7,0.6,CYAN)
    title(s,t("万亿级 IP 经济，正在向链上迁移","A Trillion-Dollar IP Economy, Moving On-chain"),y=0.95,size=30)
    txt(s,[[R(t("创作者经济持续高速增长，而「粉丝情绪可定价、可交易」是其中尚未被满足的最大缺口。",
        "The creator economy keeps surging; making fan emotion priceable & tradable is its biggest unmet gap."),MUTED,False,14)]],0.72,1.7,11.8,0.4)
    ccx,ccy=3.55,4.5
    for d,lc,fc in [(4.4,VIOLET,PANEL2),(3.0,MAGENTA,PANEL),(1.7,CYAN,PANEL2)]:
        o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(ccx-d/2),Inches(ccy-d/2),Inches(d),Inches(d))
        o.fill.solid(); o.fill.fore_color.rgb=C(fc); o.line.color.rgb=C(lc); o.line.width=Pt(2.0); shadow(o,blur=18,dist=0,alpha=35,color=lc)
    txt(s,[[R("TAM",VIOLET,True,12,DISP)],[R(t("$2,500 亿+","$250B+"),TEXT,True,17,DISP)]],ccx-1.6,2.55,3.2,0.7,align=PP_ALIGN.CENTER,leading=1.0)
    txt(s,[[R("SAM",MAGENTA,True,11,DISP)],[R(t("$500 亿","$50B"),TEXT,True,15,DISP)]],ccx-1.6,3.42,3.2,0.65,align=PP_ALIGN.CENTER,leading=1.0)
    txt(s,[[R("SOM",CYAN,True,11,DISP)],[R(t("$15 亿","$1.5B"),TEXT,True,15,DISP)]],ccx-1.0,4.2,2.0,0.6,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
    defs=[(t("TAM 全球创作者经济","TAM · Global creator economy"),t("全球内容创作者 / 网红经济年规模","Annual global creator / influencer economy"),VIOLET),
          (t("SAM 可发币的 IP 市场","SAM · Tokenizable IP market"),t("适合资产化的网红 / KOL · 品牌 IP","Influencer / KOL / brand IP fit to tokenize"),MAGENTA),
          (t("SOM 3 年可获取","SOM · Obtainable in 3 yrs"),t("OneIP 在发币+交易+社交的可触达份额","OneIP's reachable launch+trade+social share"),CYAN)]
    dy=2.15
    for h,b,col in defs:
        card(s,7.0,dy,5.6,0.92,fill=PANEL,radius=0.1)
        dot(s,7.28,dy+0.36,0.2,col,glow=True)
        txt(s,[[R(h,TEXT,True,14)]],7.65,dy+0.16,4.8,0.4)
        txt(s,[[R(b,MUTED,False,11)]],7.65,dy+0.54,4.85,0.35)
        dy+=1.04
    txt(s,[[R(t("增长驱动：","Drivers: "),CYAN,True,12.5),
            R(t("创作者经济年增 ~20% · Web3 粉丝代币需求爆发 · 短视频带来海量 KOL 供给 · 粉丝愿为「归属感」付费",
                "~20% creator-economy CAGR · surging fan-token demand · short-video KOL supply · fans pay for belonging"),TEXT,False,11.5)]],7.0,5.45,5.7,0.9,leading=1.3)
    txt(s,[[R(t("* 市场数字为公开资料估算，用于规模示意。","* Market figures are public-data estimates, for scale illustration."),DIM,False,10)]],7.0,6.5,5.6,0.3)
    pageno(s)

    # ---- S11 BUSINESS MODEL ----
    s=slide("bg_content.png")
    eyebrow(s,t("BUSINESS MODEL · 商业化","BUSINESS MODEL"),0.7,0.6,CYAN)
    title(s,t("多元收入结构，规模化的现金流","Diversified Revenue, Scalable Cash Flow"),y=0.95)
    streams=[(t("交易手续费","Trading fees"),t("oneIP.io 核心收入，随交易量线性放大","oneIP.io core revenue, scales with volume"),MAGENTA),
             (t("AI 订阅 (SaaS)","AI subscriptions (SaaS)"),t("oneIP.ai 月费 / 年费，稳定经常性收入","oneIP.ai recurring monthly / annual revenue"),CYAN),
             (t("发行 & 授权抽成","Launch & licensing"),t("代币发行服务费 + IP 授权交易分成","Launch fees + IP licensing share"),VIOLET),
             (t("增值 & 企业定制","Premium & enterprise"),t("终端订阅、做市分成、品牌定制 IP","Terminal subs, market-making, custom IP"),CORAL)]
    sy=2.2
    for h,b,col in streams:
        card(s,0.72,sy,5.7,1.0,fill=PANEL,radius=0.1)
        dot(s,1.02,sy+0.34,0.32,PANEL2); dot(s,1.12,sy+0.44,0.13,col,glow=True)
        txt(s,[[R(h,TEXT,True,15)]],1.55,sy+0.16,4.7,0.4)
        txt(s,[[R(b,MUTED,False,11)]],1.55,sy+0.55,4.75,0.4)
        sy+=1.12
    cd=CategoryChartData(); cd.categories=["Y1","Y2","Y3"]
    cd.add_series(t("营收预测 (万美金)","Revenue (US$ M)"), ((120,480,1500) if LANG=="zh" else (1.2,4.8,15)))
    gframe=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(6.75),Inches(2.2),Inches(5.85),Inches(4.3),cd)
    chart=gframe.chart; chart.has_legend=False; chart.has_title=True
    chart.chart_title.text_frame.text=t("营收增长预测（示意）","Revenue Growth (illustrative)")
    ctf=chart.chart_title.text_frame.paragraphs[0].runs[0].font; ctf.size=Pt(13); ctf.color.rgb=C(MUTED); ctf.bold=True; ctf.name=BODY
    plot=chart.plots[0]; plot.gap_width=90
    plot.series[0].format.fill.solid(); plot.series[0].format.fill.fore_color.rgb=C(VIOLET)
    plot.has_data_labels=True; dl=plot.data_labels; dl.font.size=Pt(12); dl.font.bold=True; dl.font.color.rgb=C(TEXT); dl.font.name=DISP
    if LANG=="en": dl.number_format='"$"0.0"M"'; dl.number_format_is_linked=False
    for ax in (chart.category_axis,chart.value_axis):
        ax.tick_labels.font.size=Pt(11); ax.tick_labels.font.color.rgb=C(MUTED); ax.tick_labels.font.name=DISP
    chart.value_axis.has_major_gridlines=True; chart.value_axis.major_gridlines.format.line.color.rgb=C(LINE)
    chart.value_axis.major_gridlines.format.line.width=Pt(0.5); chart.value_axis.visible=False
    pageno(s)

    # ---- S12 WEB4 ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE BIG PICTURE · Web4 增长愿景","THE BIG PICTURE · WEB4 VISION"),0.7,0.6,CYAN,w=12)
    title(s,t("Web4 时代：粉丝即资产，聚合亿级 IP 经济","Web4: Fans as Assets, an IP Economy at Scale"),y=0.95,size=29)
    txt(s,[[R(t("当平台聚合足够多的品牌、网红与粉丝数据，规模效应将驱动收入指数级增长。",
        "Once enough brands, influencers and fan data aggregate, scale effects drive exponential revenue."),MUTED,False,14)]],0.72,1.66,11.9,0.4)
    assume=[(t("入驻品牌 & 供应商","Brands & suppliers"),"30,000+",VIOLET),
            (t("网红 / 创作者 IP","Influencer / creator IP"),"200,000+",MAGENTA),
            (t("覆盖粉丝数据","Fans reached"),t("3 亿+","300M+"),PINK),
            (t("年交易额 GMV","Annual GMV"),t("$30 亿+","$3B+"),CYAN)]
    ax,aw,ag=0.6,2.85,0.24
    for lab,val,col in assume:
        card(s,ax,2.1,aw,1.5,fill=PANEL,radius=0.1)
        txt(s,[[R(lab,MUTED,False,12)]],ax+0.28,2.32,aw-0.5,0.35)
        txt(s,[[R(val,col,True,26,DISP)]],ax+0.28,2.72,aw-0.5,0.6)
        txt(s,[[R(t("2028E 规模假设","2028E assumption"),DIM,False,10)]],ax+0.28,3.3,aw-0.5,0.25)
        ax+=aw+ag
    wd=CategoryChartData(); wd.categories=["2026","2027","2028"]
    wd.add_series(t("平台年收入 (万美金)","Annual revenue (US$ M)"), ((300,2500,15000) if LANG=="zh" else (3,25,150)))
    wf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(0.7),Inches(4.05),Inches(6.9),Inches(2.65),wd)
    wc=wf.chart; wc.has_legend=False; wc.has_title=True
    wc.chart_title.text_frame.text=t("平台年收入预测（规模化情景）","Annual Revenue (scale scenario)")
    _f=wc.chart_title.text_frame.paragraphs[0].runs[0].font; _f.size=Pt(13); _f.color.rgb=C(MUTED); _f.bold=True; _f.name=BODY
    wp=wc.plots[0]; wp.gap_width=80
    for pt,cc in zip(wp.series[0].points,[VIOLET,MAGENTA,CYAN]):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb=C(cc)
    wp.has_data_labels=True; _dl=wp.data_labels
    _dl.number_format=('#,##0"万"' if LANG=="zh" else '"$"#,##0"M"'); _dl.number_format_is_linked=False
    _dl.font.size=Pt(12); _dl.font.bold=True; _dl.font.color.rgb=C(TEXT); _dl.font.name=DISP
    for axx in (wc.category_axis,wc.value_axis):
        axx.tick_labels.font.size=Pt(11); axx.tick_labels.font.color.rgb=C(MUTED); axx.tick_labels.font.name=DISP
    wc.value_axis.visible=False; wc.value_axis.has_major_gridlines=True
    wc.value_axis.major_gridlines.format.line.color.rgb=C(LINE); wc.value_axis.major_gridlines.format.line.width=Pt(0.5)
    card(s,7.95,3.95,4.65,2.85,fill=PANEL2,radius=0.09,line=VIOLET,line_w=1.4)
    txt(s,[[R(t("2028E 平台年收入","2028E platform revenue"),MUTED,False,11.5)]],8.3,4.06,4.15,0.35)
    txt(s,[[R(t("$1.5 亿","$150M"),CYAN,True,25,DISP)]],8.3,4.32,4.15,0.55)
    txt(s,[[R(t("估值 = 年收入 × PE 40 倍","Valuation = revenue × 40× PE"),DIM,False,11)]],8.3,4.94,4.15,0.32)
    txt(s,[[R(t("2028E 预测公司市值","2028E company valuation"),MUTED,False,11.5)]],8.3,5.3,4.15,0.35)
    txt(s,[[R(t("$60 亿","$6B"),PINK,True,37,DISP)]],8.3,5.54,4.15,0.75)
    txt(s,[[R(t("* 按 40× PE 测算，规模化情景下的愿景假设，非业绩承诺。","* 40× PE estimate; vision assumption under a scale scenario, not a forecast."),DIM,False,9.5)]],8.3,6.36,4.2,0.35,leading=1.1)
    pageno(s)

    # ---- S13 ROADMAP ----
    s=slide("bg_content.png")
    eyebrow(s,t("ROADMAP · 里程碑","ROADMAP"),0.7,0.6,CYAN)
    title(s,t("18 个月跑通闭环，3 年走向盈利","Loop in 18 Months, Profit in 3 Years"),y=0.95)
    phases=[("2026 Q3",t("平台 MVP 上线","Platform MVP"),[t("oneIP.io 发行 + 交易终端","oneIP.io launch + terminal"),t("首批 KOL / 网红入驻","First KOLs onboard"),t("钱包 · 行情 · 跟单","Wallet · charts · copy-trade")],CYAN),
            ("2026 Q4",t("AI 造 IP 上线","AI IP goes live"),[t("oneIP.ai 内测","oneIP.ai beta"),t("AI 数字分身 / 内容","AI avatar / content"),t("交易量与用户起量","Volume & users ramp")],VIOLET),
            ("2027 H1",t("规模化增长","Scale-up"),[t("品牌 / MCN 合作","Brand / MCN deals"),t("粉丝社交共创功能","Fan social & co-creation"),t("用户与营收放大","Users & revenue grow")],MAGENTA),
            ("2027 H2",t("出海与盈亏平衡","Global & break-even"),[t("多链 + 海外市场","Multi-chain + overseas"),t("做市 / 流动性完善","Market-making / liquidity"),t("实现盈亏平衡","Reach break-even")],PINK),
            ("2028",t("生态开放与盈利","Ecosystem & profit"),[t("开放生态 / API","Open ecosystem / API"),t("IP 资产二级市场","IP secondary market"),t("持续盈利增长","Sustained profit")],CORAL)]
    n=len(phases); cw=2.15; gap=(12.13-cw*n)/(n-1); x0=0.6; lineY=3.05
    ln=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x0+cw/2),Inches(lineY-0.01),Inches((cw+gap)*(n-1)),Inches(0.02))
    ln.fill.solid(); ln.fill.fore_color.rgb=C(LINE); ln.line.fill.background()
    for i,(q,h,items,col) in enumerate(phases):
        x=x0+i*(cw+gap); cxn=x+cw/2
        txt(s,[[R(q,col,True,15,DISP)]],x,2.45,cw,0.4,align=PP_ALIGN.CENTER)
        dot(s,cxn-0.13,lineY-0.13,0.26,col,glow=True)
        card(s,x,3.45,cw,2.95,fill=PANEL,radius=0.09)
        txt(s,[[R(h,TEXT,True,13.5)]],x+0.25,3.68,cw-0.45,0.75,leading=1.05)
        iy=4.55
        for it in items:
            dot(s,x+0.27,iy+0.08,0.1,col)
            txt(s,[[R(it,MUTED,False,10.5)]],x+0.5,iy-0.02,cw-0.65,0.6,leading=1.1)
            iy+=0.62
    pageno(s)

    # ---- S14 TEAM ----
    s=slide("bg_content.png")
    eyebrow(s,t("TEAM · 团队","TEAM"),0.7,0.6,CYAN)
    title(s,t("懂内容、懂链、懂增长的操盘团队","A Team That Knows Content, Chain & Growth"),y=0.95,size=30)
    txt(s,[[R(t("核心成员来自内容 / 社交生态与 Web3 一线，兼具流量获取与链上工程能力。",
        "Core members come from content/social and frontline Web3 — blending audience growth with on-chain engineering."),MUTED,False,14)]],0.72,1.7,11.8,0.4)
    members=[("BK","Benny Kau","Founder",VIOLET,MAGENTA),
             ("AY","Andy Yeo","Co-Founder",CYAN,VIOLET),
             ("JK","Joyce Kau","CEO",MAGENTA,PINK),
             ("SL","Sam Lee","COO",CORAL,PINK),
             ("FC","Fred Chong",t("Webtvasia CEO / MCN 顾问","Webtvasia CEO / MCN Advisor"),PINK,VIOLET)]
    mwid=2.18; mg=(12.13-mwid*5)/4; mx=0.6
    for i,(ini,name,role,c1,c2) in enumerate(members):
        x=mx+i*(mwid+mg)
        card(s,x,2.45,mwid,3.5,fill=PANEL,radius=0.09)
        av=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x+mwid/2-0.55),Inches(2.85),Inches(1.1),Inches(1.1))
        grad_fill(av,[(0,c1),(100,c2)],angle=55); av.line.fill.background(); shadow(av,blur=16,dist=0,alpha=40,color=c1)
        txt(s,[[R(ini,"FFFFFF",True,19,DISP)]],x+mwid/2-0.55,2.85,1.1,1.1,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,[[R(name,TEXT,True,15)]],x+0.08,4.25,mwid-0.16,0.4,align=PP_ALIGN.CENTER)
        txt(s,[[R(role,MUTED,False,10.5)]],x+0.12,4.68,mwid-0.24,0.9,align=PP_ALIGN.CENTER,leading=1.15)
    txt(s,[[R(t("完整团队简历与顾问 / 资源方名单可应投资人要求提供。",
        "Full team bios and advisor / partner list available on request."),DIM,False,11)]],0.72,6.2,11.8,0.3)
    pageno(s)

    # ---- S15 THE ASK ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE ASK · 融资计划","THE ASK"),0.7,0.62,CYAN)
    title(s,t("融资 200 万美金 · 让出 20% 原始股","Raising US$2M · for 20% Equity"),y=0.97)
    txt(s,[[R(t("投后估值 1,000 万美金，资金 18 个月跑通双引擎闭环并完成首轮规模化。",
        "US$10M post-money; runway of 18 months to prove the dual-engine loop and reach first-stage scale."),MUTED,False,15)]],0.72,1.72,11.8,0.45)
    GREEN="34D399"
    use=[(t("运营","Operations"),20,VIOLET),(t("开发","Product"),20,CYAN),(t("市场推广","Marketing"),30,MAGENTA),(t("网红·品牌·供应商","KOL/brand/supplier"),20,CORAL),(t("储备资金","Reserve"),10,GREEN)]
    dd=CategoryChartData(); dd.categories=[u[0] for u in use]; dd.add_series(t("资金用途","Use of funds"),tuple(u[1] for u in use))
    dframe=s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,Inches(0.75),Inches(2.35),Inches(4.7),Inches(4.5),dd)
    dchart=dframe.chart; dchart.has_legend=False; dchart.has_title=False
    for pt,(_,_,col) in zip(dchart.plots[0].series[0].points,use):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb=C(col); pt.format.line.color.rgb=C(BG); pt.format.line.width=Pt(2.5)
    dchart.plots[0].has_data_labels=True; ddl=dchart.plots[0].data_labels
    ddl.number_format='0"%"'; ddl.number_format_is_linked=False
    ddl.font.size=Pt(12); ddl.font.bold=True; ddl.font.color.rgb=C("FFFFFF"); ddl.font.name=DISP
    txt(s,[[R(t("200万","US$2M"),TEXT,True,22,DISP)],[R(t("美金","total"),MUTED,False,12)]],1.95,4.15,2.3,0.9,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
    funds=[("20%",t("运营","Operations"),t("团队、多国办公室、合规与日常运转","Team, multi-country offices, compliance & ops"),VIOLET,t("40 万美金","US$400K")),
           ("20%",t("开发","Product"),t("oneIP.ai 模型 + oneIP.io 终端迭代","oneIP.ai models + oneIP.io terminal"),CYAN,t("40 万美金","US$400K")),
           ("30%",t("市场推广","Marketing"),t("用户增长、品牌曝光、社区运营","User growth, brand & community"),MAGENTA,t("60 万美金","US$600K")),
           ("20%",t("网红·品牌·供应商","KOL · brand · supplier"),t("KOL 签约、品牌对接、供应链资源","KOL signing, brand & supply deals"),CORAL,t("40 万美金","US$400K")),
           ("10%",t("储备资金","Reserve"),t("机动资金，应对不确定性与机会","Buffer for contingencies & opportunities"),GREEN,t("20 万美金","US$200K"))]
    fx,fy=5.85,2.3
    for pct,h,b,col,amt in funds:
        card(s,fx,fy,6.75,0.78,fill=PANEL,radius=0.12)
        txt(s,[[R(pct,col,True,21,DISP)]],fx+0.32,fy+0.19,1.2,0.4,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,[[R(h,TEXT,True,14.5)]],fx+1.7,fy+0.1,3.6,0.34)
        txt(s,[[R(b,MUTED,False,10.5)]],fx+1.7,fy+0.42,4.2,0.3)
        txt(s,[[R(amt,col,True,12.5)]],fx+5.25,fy+0.19,1.35,0.4,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)
        fy+=0.86
    pageno(s)

    # ---- S16 VISION ----
    s=slide("bg_cover.png")
    eyebrow(s,t("THE VISION · 愿景","THE VISION"),0.75,0.95,CYAN)
    if LANG=="zh":
        txt(s,[[R("让每一个独一无二的 IP，",TEXT,True,36)],[R("都拥有自己的链上经济。",TEXT,True,36)]],0.78,1.7,8.5,2.0,leading=1.18,sa=4)
        txt(s,[[R("AI 取代了可被复制的一切；",MUTED,False,16)],[R("OneIP 让不可复制的「人」与「粉丝情绪」，成为真正的资产。",MUTED,False,16)]],0.78,3.65,8.4,1.2,leading=1.3,sa=3)
    else:
        txt(s,[[R("Give every unique IP",TEXT,True,36)],[R("its own on-chain economy.",TEXT,True,36)]],0.78,1.7,8.8,2.0,leading=1.14,sa=4)
        txt(s,[[R("AI replaced everything copyable;",MUTED,False,16)],[R("OneIP turns the uncopyable — people & fan emotion — into real assets.",MUTED,False,15)]],0.78,3.65,8.5,1.2,leading=1.3,sa=3)
    logo_lockup(s,0.75,4.85,1.05,46)
    txt(s,[[R(t("网红 IP 发币 · 连接粉丝 · 价值共创","Launch IP · Connect Fans · Co-create Value"),CYAN,True,15)]],0.78,6.05,8,0.4)
    txt(s,[[R("oneIP.io",PINK,True,14,DISP),R("　/　",DIM,False,14,DISP),R("oneIP.ai",CYAN,True,14,DISP),
            R("　　bossses001@gmail.com",MUTED,False,13,DISP)]],0.78,6.5,10,0.4)
    mascot(s,9.55,1.95,3.5)

# =================================================================== RUN
for _lang,_body,_out in [("zh","Microsoft YaHei","OneIP_融资计划书.pptx"),
                         ("en","Segoe UI","OneIP_Pitch_Deck_EN.pptx")]:
    LANG=_lang; BODY=_body
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    build()
    prs.save(_out)
    print("SAVED", _out, "slides:", len(prs.slides._sldIdLst))
