# -*- coding: utf-8 -*-
# OneIP 精简故事版 — 9 页（story / future / data / trends）
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image
import os, math
os.chdir(os.path.dirname(os.path.abspath(__file__)))

BG="080510"; PANEL="150E26"; PANEL2="1B1232"; TEXT="F4EFFC"; MUTED="A99FC0"; DIM="7A6F96"
CYAN="22D3EE"; VIOLET="8B5CF6"; PURPLE="A855F7"; MAGENTA="EC4899"; PINK="F472B6"; CORAL="FB7185"; LINE="2A2140"; GREEN="34D399"
GRAD=[CYAN,VIOLET,MAGENTA,PINK,CORAL]; DISP="Space Grotesk"
LANG="zh"; BODY="Microsoft YaHei"; prs=None; TOTAL=9
def t(zh,en): return zh if LANG=="zh" else en
def C(h): return RGBColor.from_string(h)

def slide(bg="bg_content.png"):
    s=prs.slides.add_slide(prs.slide_layouts[6]); s.shapes.add_picture(bg,0,0,width=prs.slide_width,height=prs.slide_height); return s
def shadow(shape,blur=20,dist=7,angle=90,color="000000",alpha=58):
    sp=shape._element.spPr; ex=sp.find(qn('a:effectLst'))
    if ex is not None: sp.remove(ex)
    sp.append(parse_xml('<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{int(blur*12700)}" dist="{int(dist*12700)}" dir="{int(angle*60000)}" rotWithShape="0">'
        f'<a:srgbClr val="{color}"><a:alpha val="{int(alpha*1000)}"/></a:srgbClr></a:outerShdw></a:effectLst>'))
def grad_fill(shape,stops,angle=45):
    sp=shape._element.spPr
    for tag in ('a:solidFill','a:noFill','a:gradFill','a:blipFill','a:pattFill'):
        e=sp.find(qn(tag))
        if e is not None: sp.remove(e)
    gs="".join(f'<a:gs pos="{int(p*1000)}"><a:srgbClr val="{c}"/></a:gs>' for p,c in stops)
    xml=('<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" rotWithShape="1">'
         f'<a:gsLst>{gs}</a:gsLst><a:lin ang="{int(angle*60000)}" scaled="1"/></a:gradFill>')
    ln=sp.find(qn('a:ln'))
    if ln is not None: ln.addprevious(parse_xml(xml))
    else: sp.append(parse_xml(xml))
def card(s,x,y,w,h,fill=PANEL,radius=0.085,line=LINE,line_w=1.0,sh=True):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    shp.adjustments[0]=radius; shp.fill.solid(); shp.fill.fore_color.rgb=C(fill)
    if line: shp.line.color.rgb=C(line); shp.line.width=Pt(line_w)
    else: shp.line.fill.background()
    if sh: shadow(shp)
    return shp
def dot(s,x,y,d,fill,glow=False):
    shp=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb=C(fill); shp.line.fill.background()
    if glow: shadow(shp,blur=22,dist=0,alpha=45,color=fill)
    return shp
def txt(s,runs,x,y,w,h,size=16,color=TEXT,bold=False,font=None,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,spc=None,leading=None,sa=None,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf,m,0)
    if isinstance(runs,str): runs=[[(runs,color,bold,size,font)]]
    elif runs and isinstance(runs[0],tuple): runs=[runs]
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.alignment=align
        if leading: p.line_spacing=leading
        if sa is not None: p.space_after=Pt(sa)
        p.space_before=Pt(0)
        for (tt,col,bd,sz,fn) in para:
            fn_use=fn if fn else BODY
            r=p.add_run(); r.text=tt; r.font.size=Pt(sz); r.font.bold=bd; r.font.color.rgb=C(col); r.font.name=fn_use
            rPr=r._r.get_or_add_rPr(); ea=rPr.find(qn('a:ea'))
            if ea is None: ea=rPr.makeelement(qn('a:ea'),{}); rPr.append(ea)
            ea.set('typeface',BODY)
            if spc is not None: rPr.set('spc',str(int(spc*100)))
    return tb
def R(tt,col=TEXT,bd=False,sz=16,fn=None): return (tt,col,bd,sz,fn)
def eyebrow(s,tt,x,y,color=MAGENTA,w=12): txt(s,[[R(tt,color,True,12.5,DISP)]],x,y,w,0.3,spc=2.4)
def title(s,tt,x=0.7,y=0.78,size=33,w=11.8,color=TEXT): txt(s,[[R(tt,color,True,size)]],x,y,w,0.95)
mw,mh=Image.open("onee.png").size; ASP=mw/mh
def mascot(s,x,y,w): s.shapes.add_picture("onee.png",Inches(x),Inches(y),width=Inches(w))
def logo_lockup(s,x,y,h,tsize,text="neIP.io",gap=0.06,w=7.5):
    s.shapes.add_picture("logo_mark.png",Inches(x),Inches(y),height=Inches(h))
    runs=[]; n=len(text)
    for i,ch in enumerate(text): runs.append(R(ch,GRAD[min(int(i*len(GRAD)/max(n,1)),len(GRAD)-1)],True,tsize,DISP))
    txt(s,[runs],x+h+gap,y,w,h,anchor=MSO_ANCHOR.MIDDLE)
def pageno(s):
    idx=len(prs.slides._sldIdLst)
    s.shapes.add_picture("logo_mark.png",Inches(0.7),Inches(7.0),height=Inches(0.33))
    txt(s,[[R("oneIP.io",MUTED,True,10,DISP),R(t("  ·  精简故事版","  ·  Story Deck"),DIM,False,10)]],1.12,7.0,5,0.33,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,[[R(f"{idx:02d} / {TOTAL}",DIM,False,10,DISP)]],11.4,7.0,1.3,0.33,align=PP_ALIGN.RIGHT,anchor=MSO_ANCHOR.MIDDLE)

def build():
    # ---- 1 COVER ----
    s=slide("bg_cover.png")
    eyebrow(s,t("STORY DECK · 精简版 9 页","STORY DECK · 9 PAGES"),0.75,0.62,CYAN)
    logo_lockup(s,0.7,1.12,1.35,66)
    txt(s,[[R("OneIP",CYAN,True,33,DISP),R(" + ",DIM,False,33,DISP),R("OPC",PINK,True,33,DISP)]],0.78,2.5,8,0.65)
    txt(s,[[R("OPC = One Personal Company",TEXT,True,14,DISP),R(t("　一人即一家公司","  ·  you are the company"),MUTED,False,14)]],0.78,3.22,8.8,0.38)
    if LANG=="zh":
        txt(s,[[R("品牌 IP 一键发币",PINK,True,18),R(" · ",DIM,False,18),R("链接粉丝",CYAN,True,18),R(" · ",DIM,False,18),R("价值共创",VIOLET,True,18)]],0.78,3.78,8.7,0.45)
    else:
        txt(s,[[R("Launch Brand IP",PINK,True,18,DISP),R(" · ",DIM,False,18,DISP),R("Connect Fans",CYAN,True,18,DISP),R(" · ",DIM,False,18,DISP),R("Co-create Value",VIOLET,True,18,DISP)]],0.78,3.78,8.9,0.45)
    chips=[(t("融资","Raise"),t("200 万美金","US$2M"),CYAN),(t("让出原始股","Equity"),"20%",PINK),(t("投后估值","Post-money"),t("1,000 万美金","US$10M"),VIOLET)]
    cx=0.78
    for lab,val,col in chips:
        card(s,cx,5.6,2.55,0.95,fill=PANEL2,radius=0.22,line=LINE,line_w=1)
        txt(s,[[R(lab,MUTED,False,11)]],cx+0.28,5.77,2.05,0.3); txt(s,[[R(val,col,True,19,DISP)]],cx+0.28,6.07,2.05,0.4); cx+=2.83
    mascot(s,9.5,1.5,3.5)

    # ---- 2 THE SHIFT ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE SHIFT · 时代变了","THE SHIFT"),0.7,0.75,CYAN)
    title(s,t("AI 正在吞噬一切可被复制的东西","AI Is Eating Everything Copyable"),y=1.2,size=34)
    txt(s,[[R(t("客服、设计、文案、编程、初级岗位——正被 AI 一一接管。价值，正从「可复制的劳动」迁移到「不可复制的人」。",
        "Support, design, copy, coding, entry-level roles — taken over by AI. Value migrates from copyable labor to the uncopyable person."),MUTED,False,15.5)]],0.72,2.15,11.8,0.7,leading=1.3)
    big=[(t("2030 AI 全球经济","2030 global AI economy"),"$4.9T",CYAN),
         (t("可自动化岗位","Automatable jobs"),"35%+",MAGENTA),
         (t("创作者经济规模","Creator economy"),"$500B",PINK)]
    bx=0.72
    for lab,val,col in big:
        card(s,bx,3.3,3.85,2.4,fill=PANEL,radius=0.09)
        txt(s,[[R(val,col,True,46,DISP)]],bx+0.4,3.75,3.2,1.0)
        txt(s,[[R(lab,TEXT,True,15)]],bx+0.4,4.95,3.2,0.5,leading=1.1)
        bx+=4.13
    txt(s,[[R(t("→ 稀缺性，成为唯一的溢价来源。","→ Scarcity becomes the only premium."),VIOLET,True,15)]],0.72,6.0,11,0.4)
    pageno(s)

    # ---- 3 MEGATREND DATA ----
    s=slide("bg_content.png")
    eyebrow(s,t("MEGATREND · 趋势数据（行业预测）","MEGATREND · INDUSTRY FORECAST"),0.7,0.58,CYAN)
    title(s,t("AI 越强，个人 IP 越值钱","The Stronger AI Gets, the More Personal IP Is Worth"),y=0.9,size=28)
    txt(s,[[R(t("机构预测（IDC · Stanford HAI · Goldman Sachs · WEF · McKinsey）：AI 经济狂飙、岗位重构，个人品牌成为出口。",
        "Institutional forecasts (IDC · Stanford HAI · Goldman Sachs · WEF · McKinsey): AI economy soars, jobs restructure, personal brand becomes the exit."),MUTED,False,12.5)]],0.72,1.54,12.1,0.4)
    yrs=["2027","2028","2029","2030"]; yx=[5.0,6.85,8.7,10.55]; yw=1.75
    txt(s,[[R(t("指标 · 行业预测","Metric · forecast"),DIM,True,10.5)]],0.85,1.98,3.9,0.3,spc=0.6)
    for x,yy in zip(yx,yrs): txt(s,[[R(yy,DIM,True,11,DISP)]],x,1.98,yw,0.3,align=PP_ALIGN.CENTER)
    trends=[(t("AI 全球经济价值","Global AI economic value"),["$3.2T","$3.8T","$4.4T","$4.9T"],CYAN),
            (t("AI 市场年成长 (CAGR)","AI market CAGR"),["25%","27%","28%","30%"],VIOLET),
            (t("企业员工减少 (预测)","Workforce cut (forecast)"),["5–10%","10–15%","15–20%","20–30%"],MAGENTA),
            (t("AI 可自动化岗位比例","Automatable jobs"),["20%","25%","30%","35%+"],CORAL),
            (t("创作者经济规模","Creator economy"),["$400B","$440B","$470B","$500B+"],PINK),
            (t("AI 技能需求","AI-skill demand"),[t("高","High"),t("很高","Higher"),t("标配","Standard"),t("必备","Must-have")],PURPLE)]
    ry,rh,gap=2.26,0.54,0.05
    for lab,vals,col in trends:
        card(s,0.6,ry,12.13,rh,fill=PANEL,radius=0.16,sh=False)
        dot(s,0.85,ry+rh/2-0.05,0.11,col)
        txt(s,[[R(lab,TEXT,True,11.5)]],1.06,ry,3.75,rh,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
        for i,(x,v) in enumerate(zip(yx,vals)):
            last=(i==3)
            txt(s,[[R(v,(col if last else MUTED),last,(14 if last else 12),DISP)]],x,ry,yw,rh,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        ry+=rh+gap
    card(s,0.6,5.85,12.13,0.72,fill=PANEL2,radius=0.16,line=VIOLET,line_w=1.3)
    txt(s,[[R(t("专业人士转型自媒体（Creator Economy）：","Professionals pivot to the Creator Economy: "),CYAN,True,13),
            R(t("医生 · 律师 · 会计 · 老师 · 地产 · 保险 纷纷经营个人 IP —— 正是 OneIP 的用户。",
                "doctors · lawyers · accountants · teachers · realtors · agents all build personal IP — OneIP's users."),TEXT,False,11.5)]],0.95,5.85,11.5,0.72,anchor=MSO_ANCHOR.MIDDLE,leading=1.08)
    txt(s,[[R(t("来源：IDC / Stanford AI Index / Goldman Sachs / WEF / McKinsey（2024–2026）；「企业员工减少」为综合机构预测区间，非官方统计。",
        "Sources: IDC / Stanford AI Index / Goldman Sachs / WEF / McKinsey (2024–2026). \"Workforce cut\" is a blended forecast range, not official data."),DIM,False,8.5)]],0.72,6.64,12.2,0.26)
    pageno(s)

    # ---- 4 INSIGHT ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE INSIGHT · 不可被取代","THE INSIGHT"),0.7,0.62,CYAN)
    title(s,t("唯独「人」和「粉丝」，AI 无法取代","Only People & Fans Can't Be Replaced"),y=0.97)
    left=[(t("个人 IP / 自媒体","Personal IP / Creators"),t("AI 只能辅助，无法成为「那个人」。个性、故事、信任——天然独一无二。","AI can assist, but can't BE that person. Personality, story, trust — inherently unique."),VIOLET),
          (t("粉丝数据 / 粉丝情绪","Fan Data / Fan Emotion"),t("真实的关注与情感无法被生成。AI 是放大器，不是替代品。","Real attention & emotion can't be generated. AI amplifies — it doesn't replace."),MAGENTA)]
    ly=2.05
    for h,b,col in left:
        card(s,0.72,ly,6.5,1.95,fill=PANEL,radius=0.08); dot(s,1.02,ly+0.32,0.22,col,glow=True)
        txt(s,[[R(h,TEXT,True,19)]],1.45,ly+0.28,5.5,0.5); txt(s,[[R(b,MUTED,False,13)]],1.02,ly+0.95,5.95,0.9,leading=1.3); ly+=2.15
    card(s,7.55,2.05,5.05,4.05,fill=PANEL2,radius=0.08,line=VIOLET,line_w=1.2)
    txt(s,[[R(t("专业人士转型自媒体","Pros pivot to creators"),CYAN,True,17)]],7.95,2.5,4.3,0.4)
    txt(s,[[R(t("医生 · 律师 · 会计 · 老师 · 地产 · 保险，都在建立自己的个人品牌。","Doctors, lawyers, accountants, teachers, realtors, agents — all building personal brands."),MUTED,False,13)]],7.95,2.95,4.35,1.0,leading=1.3)
    txt(s,[[R(t("结论","Takeaway"),PINK,True,14)]],7.95,4.15,4.3,0.35)
    txt(s,[[R(t("个人品牌需求，指数级爆发。","Demand for personal brands is exploding."),TEXT,True,18)]],7.95,4.5,4.35,0.9,leading=1.15)
    logo_lockup(s,7.95,5.5,0.72,30,w=4)
    pageno(s)

    # ---- 5 ONEIP = ANSWER ----
    s=slide("bg_content.png")
    eyebrow(s,t("THE ANSWER · OneIP = OPC","THE ANSWER · OneIP = OPC"),0.7,0.6,CYAN)
    title(s,t("OneIP：让每个人，成为一家公司","OneIP: Turn Everyone Into a Company"),y=0.95,size=30)
    txt(s,[[R(t("一个平台，把个人品牌变成可发行、可交易、可变现的资产——OPC，一人即一家公司。",
        "One platform that turns personal brand into an issuable, tradable, monetizable asset — OPC, you are the company."),MUTED,False,14)]],0.72,1.66,12,0.4)
    eng=[("oneIP.ai",t("造 IP · 供给端","Create IP · Supply"),CYAN,[t("AI 数字分身 · 内容批量生产","AI avatar · content at scale"),t("标准化 IP 资产包","Standardized IP packs"),t("持续供给，越用越强","Compounding supply")]),
         ("oneIP.io",t("卖 IP · 流通端","Trade IP · Liquidity"),PINK,[t("一键发币 · 交易终端","1-click launch · terminal"),t("链接粉丝 · 价值共创","Connect fans · co-create"),t("手续费 + 流动性变现","Fees + liquidity")])]
    ex=0.72
    for name,tag,col,items in eng:
        card(s,ex,2.25,5.62,3.9,fill=PANEL,radius=0.07)
        txt(s,[[R(name,col,True,30,DISP)]],ex+0.45,2.55,5.0,0.6); txt(s,[[R(tag,TEXT,True,15)]],ex+0.45,3.2,5.0,0.4)
        iy=3.85
        for it in items:
            dot(s,ex+0.48,iy+0.07,0.14,col); txt(s,[[R(it,MUTED,False,13.5)]],ex+0.78,iy-0.04,4.4,0.4); iy+=0.6
        ex+=5.94
    txt(s,[[R("↔",VIOLET,True,30,DISP)]],6.34,4.0,0.66,0.6,align=PP_ALIGN.CENTER)
    txt(s,[[R(t("造 IP × 卖 IP ＝ 有真实资产支撑的 IP 经济闭环。","Create IP × Trade IP = a real-asset-backed IP economy."),CYAN,True,14)]],0.72,6.35,12,0.4,align=PP_ALIGN.CENTER)
    pageno(s)

    # ---- 6 MARKET ----
    s=slide("bg_content.png")
    eyebrow(s,t("MARKET · 市场规模","MARKET"),0.7,0.6,CYAN)
    title(s,t("万亿级 IP 经济，正在向链上迁移","A Trillion-Dollar IP Economy, Moving On-chain"),y=0.95,size=30)
    txt(s,[[R(t("创作者经济高速增长，而「粉丝情绪可定价、可交易」是尚未被满足的最大缺口。",
        "The creator economy is surging; making fan emotion priceable & tradable is its biggest unmet gap."),MUTED,False,14)]],0.72,1.7,11.8,0.4)
    ccx,ccy=3.55,4.5
    for d,lc,fc in [(4.4,VIOLET,PANEL2),(3.0,MAGENTA,PANEL),(1.7,CYAN,PANEL2)]:
        o=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(ccx-d/2),Inches(ccy-d/2),Inches(d),Inches(d))
        o.fill.solid(); o.fill.fore_color.rgb=C(fc); o.line.color.rgb=C(lc); o.line.width=Pt(2.0); shadow(o,blur=18,dist=0,alpha=35,color=lc)
    txt(s,[[R("TAM",VIOLET,True,12,DISP)],[R(t("$2,500 亿+","$250B+"),TEXT,True,17,DISP)]],ccx-1.6,2.55,3.2,0.7,align=PP_ALIGN.CENTER,leading=1.0)
    txt(s,[[R("SAM",MAGENTA,True,11,DISP)],[R(t("$500 亿","$50B"),TEXT,True,15,DISP)]],ccx-1.6,3.42,3.2,0.65,align=PP_ALIGN.CENTER,leading=1.0)
    txt(s,[[R("SOM",CYAN,True,11,DISP)],[R(t("$15 亿","$1.5B"),TEXT,True,15,DISP)]],ccx-1.0,4.2,2.0,0.6,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
    defs=[(t("TAM 全球创作者经济","TAM · Global creator economy"),VIOLET),(t("SAM 可发币的 IP 市场","SAM · Tokenizable IP market"),MAGENTA),(t("SOM 3 年可获取","SOM · Obtainable in 3 yrs"),CYAN)]
    dy=2.5
    for h,col in defs:
        card(s,7.0,dy,5.6,0.85,fill=PANEL,radius=0.12); dot(s,7.3,dy+0.34,0.18,col,glow=True)
        txt(s,[[R(h,TEXT,True,15)]],7.62,dy,4.8,0.85,anchor=MSO_ANCHOR.MIDDLE); dy+=0.98
    txt(s,[[R(t("增长驱动：","Drivers: "),CYAN,True,12.5),R(t("创作者经济年增 ~20% · 粉丝代币需求爆发 · 海量 KOL 供给 · 粉丝愿为归属感付费",
        "~20% creator-economy CAGR · surging fan-token demand · KOL supply · fans pay for belonging"),TEXT,False,11.5)]],7.0,5.55,5.7,0.7,leading=1.3)
    txt(s,[[R(t("* 市场数字为公开资料估算。","* Public-data estimates."),DIM,False,10)]],7.0,6.4,5.6,0.3)
    pageno(s)

    # ---- 7 FLYWHEEL ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE FLYWHEEL · 自增长飞轮","THE FLYWHEEL"),0.7,0.62,CYAN)
    title(s,t("AI 供给 × 链上流通 = 自增长飞轮","AI Supply × On-chain Liquidity = Flywheel"),y=0.97,size=30)
    txt(s,[[R(t("别人只做发币（投机、易归零）；我们用 AI 保证真实供给，用链上提供流动性与变现。",
        "Others just launch tokens (speculative). We secure real supply via AI and liquidity via chain."),MUTED,False,15)]],0.72,1.72,11.8,0.5,leading=1.25)
    cxp,cyp=6.666,4.45; RX,RY,nw,nh=3.7,1.75,3.5,0.95
    nodes=[(t("AI 持续生产 IP 资产","AI keeps producing IP assets"),CYAN,0.0),(t("链上发行 · 提供流动性","On-chain launch · liquidity"),PINK,90.0),
           (t("交易变现 · 反哺创作者","Trading monetizes · rewards creators"),MAGENTA,180.0),(t("更多 IP / 品牌入驻","More IPs / brands onboard"),VIOLET,270.0)]
    hub=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cxp-1.05),Inches(cyp-1.05),Inches(2.1),Inches(2.1))
    grad_fill(hub,[(0,CYAN),(45,VIOLET),(100,MAGENTA)],angle=60); hub.line.fill.background(); shadow(hub,blur=26,dist=0,alpha=40,color=VIOLET)
    txt(s,[[R("OneIP","FFFFFF",True,21,DISP)],[R(t("飞轮","Flywheel"),"FFFFFF",True,13)]],cxp-1.05,cyp-0.62,2.1,1.2,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.0)
    for label,col,ang in nodes:
        a=math.radians(ang); nx=cxp+RX*math.cos(a)-nw/2; ny=cyp+RY*math.sin(a)-nh/2
        card(s,nx,ny,nw,nh,fill=PANEL2,radius=0.2,line=col,line_w=1.3)
        txt(s,[[R(label,TEXT,True,13.5)]],nx+0.2,ny,nw-0.4,nh,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,leading=1.05)
    pageno(s)

    # ---- 8 WEB4 VISION + SCALE ----
    s=slide("bg_section.png")
    eyebrow(s,t("THE FUTURE · Web4 增长愿景","THE FUTURE · WEB4 VISION"),0.7,0.6,CYAN)
    title(s,t("Web4 时代：粉丝即资产，聚合亿级 IP 经济","Web4: Fans as Assets, an IP Economy at Scale"),y=0.95,size=28)
    txt(s,[[R(t("当平台聚合足够多的品牌、网红与粉丝数据，规模效应将驱动收入指数级增长。",
        "Once enough brands, influencers and fan data aggregate, scale effects drive exponential revenue."),MUTED,False,14)]],0.72,1.66,11.9,0.4)
    assume=[(t("入驻品牌 & 供应商","Brands & suppliers"),"30,000+",VIOLET),(t("网红 / 创作者 IP","Influencer / creator IP"),"200,000+",MAGENTA),
            (t("覆盖粉丝数据","Fans reached"),t("3 亿+","300M+"),PINK),(t("年交易额 GMV","Annual GMV"),t("$30 亿+","$3B+"),CYAN)]
    ax,aw,ag=0.6,2.85,0.24
    for lab,val,col in assume:
        card(s,ax,2.1,aw,1.5,fill=PANEL,radius=0.1)
        txt(s,[[R(lab,MUTED,False,12)]],ax+0.28,2.32,aw-0.5,0.35); txt(s,[[R(val,col,True,26,DISP)]],ax+0.28,2.72,aw-0.5,0.6)
        txt(s,[[R(t("2028E 规模假设","2028E assumption"),DIM,False,10)]],ax+0.28,3.3,aw-0.5,0.25); ax+=aw+ag
    wd=CategoryChartData(); wd.categories=["2026","2027","2028"]
    wd.add_series(t("平台年收入 (万美金)","Annual revenue (US$ M)"), ((300,2500,15000) if LANG=="zh" else (3,25,150)))
    wf=s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(0.7),Inches(4.05),Inches(6.9),Inches(2.65),wd)
    wc=wf.chart; wc.has_legend=False; wc.has_title=True
    wc.chart_title.text_frame.text=t("平台年收入预测（规模化情景）","Annual Revenue (scale scenario)")
    _f=wc.chart_title.text_frame.paragraphs[0].runs[0].font; _f.size=Pt(13); _f.color.rgb=C(MUTED); _f.bold=True; _f.name=BODY
    wp=wc.plots[0]; wp.gap_width=80
    for pt,cc in zip(wp.series[0].points,[VIOLET,MAGENTA,CYAN]): pt.format.fill.solid(); pt.format.fill.fore_color.rgb=C(cc)
    wp.has_data_labels=True; _dl=wp.data_labels; _dl.number_format=('#,##0"万"' if LANG=="zh" else '"$"#,##0"M"'); _dl.number_format_is_linked=False
    _dl.font.size=Pt(12); _dl.font.bold=True; _dl.font.color.rgb=C(TEXT); _dl.font.name=DISP
    for axx in (wc.category_axis,wc.value_axis): axx.tick_labels.font.size=Pt(11); axx.tick_labels.font.color.rgb=C(MUTED); axx.tick_labels.font.name=DISP
    wc.value_axis.visible=False; wc.value_axis.has_major_gridlines=True; wc.value_axis.major_gridlines.format.line.color.rgb=C(LINE); wc.value_axis.major_gridlines.format.line.width=Pt(0.5)
    card(s,7.95,3.95,4.65,2.85,fill=PANEL2,radius=0.09,line=VIOLET,line_w=1.4)
    txt(s,[[R(t("2028E 平台年收入","2028E platform revenue"),MUTED,False,11.5)]],8.3,4.06,4.15,0.35)
    txt(s,[[R(t("$1.5 亿","$150M"),CYAN,True,25,DISP)]],8.3,4.32,4.15,0.55)
    txt(s,[[R(t("估值 = 年收入 × PE 40 倍","Valuation = revenue × 40× PE"),DIM,False,11)]],8.3,4.94,4.15,0.32)
    txt(s,[[R(t("2028E 预测公司市值","2028E company valuation"),MUTED,False,11.5)]],8.3,5.3,4.15,0.35)
    txt(s,[[R(t("$60 亿","$6B"),PINK,True,37,DISP)]],8.3,5.54,4.15,0.75)
    txt(s,[[R(t("* 按 40× PE 测算，规模化情景下的愿景假设，非业绩承诺。","* 40× PE estimate; vision assumption under a scale scenario, not a forecast."),DIM,False,9.5)]],8.3,6.36,4.2,0.35,leading=1.1)
    pageno(s)

    # ---- 9 ASK + CLOSE ----
    s=slide("bg_cover.png")
    eyebrow(s,t("THE ASK · 融资 & 愿景","THE ASK · JOIN US"),0.75,0.7,CYAN)
    txt(s,[[R(t("融资 200 万美金 · 让出 20% 原始股","Raising US$2M · for 20% Equity"),TEXT,True,29)]],0.78,1.28,9,0.6)
    txt(s,[[R(t("投后估值 1,000 万美金 · 18 个月跑通双引擎闭环。","US$10M post-money · 18-month runway to prove the dual-engine loop."),MUTED,False,14)]],0.78,2.02,8.6,0.4)
    if LANG=="zh":
        txt(s,[[R("让每一个独一无二的 IP，",TEXT,True,32)],[R("都拥有自己的链上经济。",TEXT,True,32)]],0.78,2.85,8.6,1.6,leading=1.16,sa=4)
    else:
        txt(s,[[R("Give every unique IP",TEXT,True,32)],[R("its own on-chain economy.",TEXT,True,32)]],0.78,2.85,8.8,1.6,leading=1.12,sa=4)
    logo_lockup(s,0.78,4.75,1.0,44)
    txt(s,[[R("OneIP + OPC",CYAN,True,15,DISP),R(t("　·　独一无二","　·　One of a Kind"),PINK,True,15)]],0.8,5.85,8,0.4)
    txt(s,[[R("oneIP.io",PINK,True,14,DISP),R("　/　",DIM,False,14,DISP),R("oneIP.ai",CYAN,True,14,DISP),R("　　bossses001@gmail.com",MUTED,False,13,DISP)]],0.8,6.3,10,0.4)
    mascot(s,9.55,1.9,3.5)

for _lang,_body,_out in [("zh","Microsoft YaHei","OneIP_精简版_9页.pptx"),("en","Segoe UI","OneIP_Short_9p_EN.pptx")]:
    LANG=_lang; BODY=_body
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    build(); prs.save(_out); print("SAVED",_out,"slides:",len(prs.slides._sldIdLst))
